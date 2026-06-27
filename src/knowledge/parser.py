# src/knowledge/parser.py
import os
import hashlib

# 文件解析依赖库（请在 requirements.txt 中安装）
try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None


def parse_file(file_path: str) -> str:
    """
    根据文件扩展名解析文件内容，返回纯文本字符串。
    支持格式: .txt, .pdf, .docx, .pptx, .xlsx
    若格式不支持或解析失败，返回空字符串。
    """
    ext = os.path.splitext(file_path)[1].lower()
    content = ""

    try:
        if ext == ".txt":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()

        elif ext == ".pdf":
            if fitz is None:
                print("警告: PyMuPDF (fitz) 未安装，无法解析 PDF")
                return ""
            doc = fitz.open(file_path)
            for page in doc:
                content += page.get_text()
            doc.close()

        elif ext == ".docx":
            if Document is None:
                print("警告: python-docx 未安装，无法解析 .docx")
                return ""
            doc = Document(file_path)
            for para in doc.paragraphs:
                if para.text:
                    content += para.text + "\n"

        elif ext == ".pptx":
            if Presentation is None:
                print("警告: python-pptx 未安装，无法解析 .pptx")
                return ""
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        content += shape.text + "\n"

        elif ext in [".xlsx", ".xls"]:
            if load_workbook is None:
                print("警告: openpyxl 未安装，无法解析 .xlsx")
                return ""
            # 只读模式，节省内存
            wb = load_workbook(file_path, read_only=True, data_only=True)
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = " ".join([str(cell) for cell in row if cell is not None])
                    if row_text:
                        content += row_text + "\n"
            wb.close()

        else:
            # 不支持的文件格式，静默跳过
            pass

    except Exception as e:
        print(f"解析文件失败 {file_path}: {e}")
        return ""

    return content.strip()


def compute_file_hash(file_path: str) -> str:
    """
    计算文件的 MD5 哈希值，用于判断文件内容是否发生变化。
    """
    hasher = hashlib.md5()
    try:
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(65536), b""):
                hasher.update(chunk)
        return hasher.hexdigest()
    except Exception as e:
        print(f"计算文件哈希失败 {file_path}: {e}")
        return ""


def get_file_metadata(file_path: str) -> dict:
    """
    获取文件的元数据，包含：
    - mtime: 最后修改时间戳
    - size: 文件大小（字节）
    - hash: 文件 MD5 值
    """
    stat = os.stat(file_path)
    return {
        "mtime": stat.st_mtime,
        "size": stat.st_size,
        "hash": compute_file_hash(file_path),
    }